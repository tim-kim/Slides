from library.config import *
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

def create_title_slide(prs, song, artist):
    p = setup_slide_layout(prs)
    run = p.add_run()
    run.text = song
    font = run.font
    font.name = FONT_NAME
    font.size = FONT_SIZE_TITLE
    font.color.rgb = FONT_WHITE

    run = p.add_run()
    run.text = '\n' + artist
    font = run.font
    font.name = FONT_NAME
    font.size = FONT_SIZE
    font.color.rgb = FONT_WHITE

def create_lyric_slide(prs, block):
    p = setup_slide_layout(prs)
    run = p.add_run()
    run.text = '\n'.join(block).upper()
    font = run.font
    font.name = FONT_NAME
    font.size = FONT_SIZE
    font.color.rgb = FONT_WHITE

def setup_slide_layout(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = FONT_BLACK

    text_box = slide.shapes.add_textbox(LEFT, TOP, prs.slide_width, HEIGHT)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    return p
